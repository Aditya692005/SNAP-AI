"""
RAG microservice — FastAPI + LangChain LCEL + ChromaDB + Gemini 2.5 Flash
Start with: uvicorn main:app --host 0.0.0.0 --port 8000 --reload
"""

import json
import os
import re
import shutil
from datetime import datetime
from pathlib import Path
from typing import Literal

import markdown
import pandas as pd
from fpdf import FPDF
from dotenv import load_dotenv
from fastapi import FastAPI, File, HTTPException, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse
from pydantic import BaseModel

from langchain_chroma import Chroma
from langchain_community.document_loaders import CSVLoader, PyPDFLoader, TextLoader
from langchain_core.documents import Document
from langchain_core.messages import AIMessage, HumanMessage

from langchain_core.output_parsers import StrOutputParser
from langchain_core.prompts import ChatPromptTemplate, MessagesPlaceholder
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain_google_genai.chat_models import ChatGoogleGenerativeAIError
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_text_splitters import RecursiveCharacterTextSplitter

load_dotenv()

GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY")
if not GOOGLE_API_KEY:
    raise RuntimeError("GOOGLE_API_KEY is not set in .env")

# Model is configurable so you can switch to one with available quota without
# code changes (each Gemini model has its own free-tier daily request limit).
GEMINI_MODEL = os.getenv("GEMINI_MODEL", "gemini-3.0-flash")

# ── Paths ──────────────────────────────────────────────────────────────────────
CHROMA_DIR = Path("./chroma_db")
UPLOAD_DIR = Path("./uploads")
UPLOAD_DIR.mkdir(exist_ok=True)

# ── LangChain components ───────────────────────────────────────────────────────
embeddings = HuggingFaceEmbeddings(
    model_name="all-MiniLM-L6-v2",
    model_kwargs={"device": "cpu"},
)

llm = ChatGoogleGenerativeAI(
    model=GEMINI_MODEL,
    google_api_key=GOOGLE_API_KEY,
    temperature=0.2,
)

splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=150)

# ── ChromaDB vector store (persisted) ─────────────────────────────────────────
vectorstore = Chroma(
    collection_name="snap_ai_docs",
    embedding_function=embeddings,
    persist_directory=str(CHROMA_DIR),
)

# ── Per-session chat history (list of HumanMessage / AIMessage) ───────────────
session_histories: dict[str, list] = {}

# ── Prompts ────────────────────────────────────────────────────────────────────
CONDENSE_PROMPT = ChatPromptTemplate.from_messages([
    MessagesPlaceholder(variable_name="chat_history"),
    ("human", "{question}"),
    ("human",
     "Given the conversation above, rewrite the follow-up question as a "
     "standalone question that captures all necessary context."),
])

RAG_PROMPT = ChatPromptTemplate.from_messages([
    ("system",
     "You are SNAP AI, a helpful enterprise assistant. "
     "Answer the user's question using ONLY the context provided below. "
     "If the context does not contain the answer, say so honestly.\n\n"
     "Context:\n{context}"),
    MessagesPlaceholder(variable_name="chat_history"),
    ("human", "{question}"),
])

PLAIN_PROMPT = ChatPromptTemplate.from_messages([
    ("system", "You are SNAP AI, a helpful enterprise assistant."),
    MessagesPlaceholder(variable_name="chat_history"),
    ("human", "{question}"),
])

# ── FastAPI app ────────────────────────────────────────────────────────────────
app = FastAPI(title="SNAP AI RAG Service")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:5173", "http://localhost:5000"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# ── File loaders ───────────────────────────────────────────────────────────────
ALLOWED_EXTENSIONS = {".pdf", ".csv", ".txt", ".xlsx", ".xls", ".docx", ".pptx"}


def excel_engine(path: Path) -> Literal["openpyxl", "xlrd"]:
    """Pick the right pandas/Excel engine: openpyxl reads .xlsx, xlrd reads the
    legacy .xls binary format (openpyxl cannot — it expects a zip-based .xlsx)."""
    return "xlrd" if path.suffix.lower() == ".xls" else "openpyxl"


def load_xlsx(path: Path) -> list[Document]:
    xl = pd.ExcelFile(path, engine=excel_engine(path))
    docs = []
    for sheet in xl.sheet_names:
        df = xl.parse(sheet)
        text = df.to_string(index=False)
        docs.append(Document(page_content=text, metadata={"sheet": sheet}))
    return docs


def load_docx(path: Path) -> list[Document]:
    from docx import Document as DocxDocument

    doc = DocxDocument(str(path))
    parts = [p.text for p in doc.paragraphs if p.text.strip()]

    # Include table cell text, which is not part of paragraphs.
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))

    text = "\n".join(parts)
    return [Document(page_content=text)]


def load_pptx(path: Path) -> list[Document]:
    from pptx import Presentation

    prs = Presentation(str(path))
    docs = []
    for i, slide in enumerate(prs.slides, start=1):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.textframe.paragraphs:
                    line = "".join(run.text for run in paragraph.runs).strip()
                    if line:
                        parts.append(line)
            elif shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        parts.append(" | ".join(cells))
        text = "\n".join(parts)
        if text.strip():
            docs.append(Document(page_content=text, metadata={"slide": i}))
    return docs


def load_file(path: Path) -> list[Document]:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return PyPDFLoader(str(path)).load()
    if suffix == ".csv":
        return CSVLoader(str(path)).load()
    if suffix == ".txt":
        return TextLoader(str(path)).load()
    if suffix in (".xlsx", ".xls"):
        return load_xlsx(path)
    if suffix == ".docx":
        return load_docx(path)
    if suffix == ".pptx":
        return load_pptx(path)
    raise ValueError(f"Unsupported file type: {suffix}")


def index_file(path: Path, filename: str) -> int:
    """Load, chunk, and add a file on disk to the vector store. Returns the
    number of chunks indexed."""
    docs = load_file(path)
    chunks = splitter.split_documents(docs)
    for chunk in chunks:
        chunk.metadata["source"] = filename
    vectorstore.add_documents(chunks)
    return len(chunks)


# ── Helpers ────────────────────────────────────────────────────────────────────
def get_history(session_id: str) -> list:
    return session_histories.setdefault(session_id, [])


def format_docs(docs: list[Document]) -> str:
    return "\n\n---\n\n".join(d.page_content for d in docs)


# Max characters of context fed to the model when summarizing a whole document.
# Gemini 2.5 Flash has a very large context window; this keeps requests fast/cheap.
MAX_CONTEXT_CHARS = 100_000

# Phrases that signal the user wants the whole document, not just similar snippets.
WHOLE_DOC_KEYWORDS = (
    "summary", "summarize", "summarise", "overview", "tl;dr", "tldr",
    "key points", "main points", "key takeaways", "takeaways", "gist",
    "what is this document about", "what's this document about",
    "describe this document", "describe the document", "outline",
    "what does this document", "main idea", "in a nutshell",
)


def wants_full_document(question: str) -> bool:
    q = question.lower()
    return any(k in q for k in WHOLE_DOC_KEYWORDS)


# Phrases that signal the user wants a chart/table rendered, not a prose answer.
CHART_KEYWORDS = (
    "chart", "graph", "plot", "visualize", "visualise", "visualization",
    "visualisation", "bar chart", "bar graph", "pie chart", "line chart",
    "line graph", "histogram", "scatter", "diagram", "draw a", "plot a",
    "show me a", "pie of", "breakdown of",
)


def wants_chart(question: str) -> bool:
    q = question.lower()
    return any(k in q for k in CHART_KEYWORDS)


# A generate-style verb followed (allowing filler words) by a document noun
# signals the user wants a generated report/PDF — e.g. "make me a report",
# "create a detailed PDF of the findings", "draft a one-page brief".
DOC_RE = re.compile(
    r"\b(generate|create|write|draft|prepare|produce|make|build|compile|export|give me)\b"
    r".{0,40}?"
    r"\b(report|document|pdf|memo|brief|write[-\s]?up|white\s?paper|dossier)\b",
    re.IGNORECASE,
)


def wants_document(question: str) -> bool:
    return bool(DOC_RE.search(question))


def get_chunks_for_source(source: str) -> list[Document]:
    """All indexed chunks belonging to one uploaded file (insertion order)."""
    res = vectorstore._collection.get(
        where={"source": source}, include=["documents", "metadatas"]
    )
    documents = res.get("documents") or []
    metadatas = res.get("metadatas") or []
    return [
        Document(page_content=text or "", metadata=meta or {"source": source})
        for text, meta in zip(documents, metadatas)
    ]


def list_sources() -> list[str]:
    """Distinct filenames of every indexed document."""
    results = vectorstore._collection.get(include=["metadatas"])
    sources = {str(m.get("source", "")) for m in (results.get("metadatas") or [])}
    return [s for s in sources if s]


def resolve_source(source: str | None) -> str | None:
    """Pick the document to chart: the requested one, or the only one indexed.
    Returns None when ambiguous (multiple docs, none specified)."""
    if source:
        return source
    sources = list_sources()
    return sources[0] if len(sources) == 1 else None


# ── Visualization helpers ────────────────────────────────────────────────────
# Max characters of dataset preview fed to the model when building a chart spec.
MAX_VIZ_CHARS = 60_000


def load_dataframe(source: str) -> "pd.DataFrame | None":
    """Re-read the original uploaded file as a DataFrame (CSV / Excel only).

    Charts need accurate numeric data, so for tabular files we parse the source
    on disk rather than relying on the chunked/embedded text. Returns None for
    non-tabular files or if parsing fails.
    """
    safe_name = Path(source).name
    path = (UPLOAD_DIR / safe_name).resolve()
    if UPLOAD_DIR.resolve() not in path.parents or not path.is_file():
        return None
    suffix = path.suffix.lower()
    try:
        if suffix == ".csv":
            return pd.read_csv(path)
        if suffix in (".xlsx", ".xls"):
            return pd.read_excel(path, engine=excel_engine(path))  # first sheet
    except Exception:
        return None
    return None


def build_viz_context(source: str | None) -> tuple[str, bool]:
    """Return (dataset_text, is_tabular) used as context for chart generation.

    When `source` is None (e.g. the user asked for a chart but several documents
    are indexed and none is focused), aggregate text across ALL documents so a
    chart can still be produced instead of silently falling back to prose."""
    if source:
        df = load_dataframe(source)
        if df is not None:
            preview = df.head(500)
            text = "TABULAR DATA (CSV, first rows):\n" + preview.to_csv(index=False)
            return text[:MAX_VIZ_CHARS], True
        docs = get_chunks_for_source(source)
    else:
        results = vectorstore._collection.get(include=["documents", "metadatas"])
        documents = results.get("documents") or []
        metadatas = results.get("metadatas") or []
        docs = [
            Document(page_content=t or "", metadata=m or {})
            for t, m in zip(documents, metadatas)
        ]
    text = "DOCUMENT TEXT:\n" + format_docs(docs)
    return text[:MAX_VIZ_CHARS], False


VIZ_PROMPT = ChatPromptTemplate.from_messages([
    ("system",
     "You are a data-visualization assistant. Given a dataset and a user's "
     "request, produce ONE chart or table specification as STRICT JSON only — "
     "no markdown fences, no prose.\n\n"
     "JSON shape:\n"
     "{{\n"
     '  "chart_type": one of "bar","line","area","pie","doughnut","scatter","table",\n'
     '  "title": short descriptive title,\n'
     '  "labels": [category / x-axis labels as strings],   // for non-table charts\n'
     '  "datasets": [{{"label": series name, "data": [numbers]}}],  // for non-table charts\n'
     '  "table_columns": [column names],     // ONLY when chart_type is "table"\n'
     '  "table_rows": [[cell values as strings]],  // ONLY when chart_type is "table"\n'
     '  "notes": one short sentence describing what is shown\n'
     "}}\n\n"
     "Rules:\n"
     "- Pick the most suitable chart_type for the request and the data.\n"
     "- Use ONLY values present in the dataset; never invent data.\n"
     "- Aggregate/group when the request implies it (e.g. totals by category).\n"
     "- labels and each dataset's data array must be the SAME length.\n"
     "- pie/doughnut must have exactly one dataset.\n"
     "- Numbers must be plain (no commas, currency symbols, or units).\n"
     "- If the data cannot support the request, return a sensible table instead.\n"
     "- Output ONLY the JSON object."),
    ("human", "DATASET:\n{dataset}\n\nREQUEST: {instruction}"),
])


def parse_chart_json(raw: str) -> dict:
    """Parse the model's JSON output, tolerating ```json fences / stray text."""
    text = raw.strip()
    if text.startswith("```"):
        text = text.split("```", 2)[1] if "```" in text[3:] else text.strip("`")
        if text.lstrip().lower().startswith("json"):
            text = text.lstrip()[4:]
    # Fall back to the outermost {...} if there's surrounding prose.
    start, end = text.find("{"), text.rfind("}")
    if start != -1 and end != -1 and end > start:
        text = text[start:end + 1]
    return json.loads(text)


async def run_visualization(source: str | None, instruction: str) -> dict:
    dataset, is_tabular = build_viz_context(source)
    if not dataset.replace("DOCUMENT TEXT:", "").strip():
        raise HTTPException(status_code=422, detail="No document data to chart yet.")
    chain = VIZ_PROMPT | llm | StrOutputParser()
    raw = await chain.ainvoke({"dataset": dataset, "instruction": instruction})
    try:
        spec = parse_chart_json(raw)
    except (json.JSONDecodeError, ValueError):
        raise HTTPException(
            status_code=422,
            detail="Could not turn that request into a chart. Try rephrasing, "
                   "e.g. 'bar chart of total sales by region'.",
        )
    spec["source"] = source or "all documents"
    spec["is_tabular"] = is_tabular
    return spec


# ── Dashboard metric extraction ──────────────────────────────────────────────
# Department-specific metric catalog the dashboard understands. Maps each
# canonical metric key to the department it belongs to. The prompt maps synonyms
# onto these keys. Extraction runs for ALL of these regardless of what the user
# has chosen to display on their dashboard.
METRIC_CATALOG = {
    # Finance
    "revenue": "finance",
    "profit": "finance",
    "expenditure": "finance",
    "cash_flow": "finance",
    # Sales
    "sales": "sales",
    "units_sold": "sales",
    "new_customers": "sales",
    "average_deal_size": "sales",
    # Marketing
    "marketing_spend": "marketing",
    "leads": "marketing",
    "conversion_rate": "marketing",
    "website_traffic": "marketing",
    # Human Resources
    "headcount": "hr",
    "attrition_rate": "hr",
    "new_hires": "hr",
    "training_cost": "hr",
    # Operations
    "production_output": "operations",
    "defect_rate": "operations",
    "inventory": "operations",
    "on_time_delivery": "operations",
}
CANONICAL_METRICS = tuple(METRIC_CATALOG.keys())

# Human-readable enumeration of allowed keys grouped by department, fed to the
# model so it knows exactly which metric keys it may emit.
_CATALOG_BY_DEPT: dict[str, list[str]] = {}
for _k, _d in METRIC_CATALOG.items():
    _CATALOG_BY_DEPT.setdefault(_d, []).append(_k)
ALLOWED_METRICS_TEXT = "\n".join(
    f"- {dept}: {', '.join(keys)}" for dept, keys in _CATALOG_BY_DEPT.items()
)

METRICS_PROMPT = ChatPromptTemplate.from_messages([
    ("system",
     "You extract department-specific business metrics from a document for a "
     "dashboard. Return STRICT JSON only: a JSON ARRAY of metric objects, "
     "nothing else.\n\n"
     "Allowed metric keys, grouped by department:\n"
     f"{ALLOWED_METRICS_TEXT}\n\n"
     "Map common synonyms onto these keys, e.g.: expenses/costs->expenditure, "
     "net income/earnings->profit, turnover/total revenue->revenue, "
     "operating cash->cash_flow, units/quantity sold->units_sold, "
     "new accounts->new_customers, ad spend/campaign budget->marketing_spend, "
     "MQLs/prospects->leads, conv. rate/close rate->conversion_rate, "
     "visits/sessions->website_traffic, employees/staff->headcount, "
     "turnover rate/churn (staff)->attrition_rate, recruits->new_hires, "
     "output/production volume->production_output, defects/scrap->defect_rate, "
     "stock on hand->inventory, OTD/on-time %->on_time_delivery.\n\n"
     "Each object:\n"
     "{{\n"
     '  "metric": one of the allowed keys above,\n'
     '  "department": the department that key belongs to,\n'
     '  "period": "YYYY" | "YYYY-MM" | "YYYY-Qn", or null if none stated,\n'
     '  "value": a plain number (no commas, currency symbols, %, or units),\n'
     '  "currency": ISO code like "USD" if known, else null,\n'
     '  "category": a breakdown label (e.g. region/department/product) or null,\n'
     '  "confidence": 0.0-1.0\n'
     "}}\n\n"
     "Rules:\n"
     "- Use ONLY numbers present in the document; never invent values.\n"
     "- Emit one object per (metric, period, category) data point you find.\n"
     "- For rates/percentages emit the plain number (e.g. 12.5 for 12.5%).\n"
     "- Ignore anything that does not map to an allowed key.\n"
     "- If the document contains no such metrics, return [].\n"
     "- Output ONLY the JSON array."),
    ("human", "DOCUMENT:\n{dataset}"),
])


def parse_json_array(raw: str) -> list:
    """Parse a JSON array from the model output, tolerating fences/stray text."""
    text = raw.strip()
    if text.startswith("```"):
        text = text.split("```", 2)[1] if "```" in text[3:] else text.strip("`")
        if text.lstrip().lower().startswith("json"):
            text = text.lstrip()[4:]
    start, end = text.find("["), text.rfind("]")
    if start != -1 and end != -1 and end > start:
        text = text[start:end + 1]
    data = json.loads(text)
    return data if isinstance(data, list) else []


def normalize_metric(raw: dict) -> dict | None:
    """Validate/normalize one extracted metric; drop anything unusable."""
    metric = str(raw.get("metric", "")).strip().lower()
    if metric not in CANONICAL_METRICS:
        return None
    value_raw = raw.get("value")
    if value_raw is None:
        return None
    try:
        value = float(value_raw)
    except (TypeError, ValueError):
        return None
    period = raw.get("period")
    category = raw.get("category")
    currency = raw.get("currency")
    try:
        confidence = float(raw.get("confidence", 0.5))
    except (TypeError, ValueError):
        confidence = 0.5
    return {
        "metric": metric,
        "department": METRIC_CATALOG[metric],
        "period": str(period) if period not in (None, "") else None,
        "value": value,
        "currency": str(currency) if currency else None,
        "category": str(category) if category else None,
        "confidence": max(0.0, min(1.0, confidence)),
    }


async def extract_metrics(source: str) -> list[dict]:
    """Pull canonical financial metrics from one uploaded document."""
    dataset, _ = build_viz_context(source)
    if not dataset.strip():
        return []
    chain = METRICS_PROMPT | llm | StrOutputParser()
    raw = await chain.ainvoke({"dataset": dataset})
    try:
        items = parse_json_array(raw)
    except (json.JSONDecodeError, ValueError):
        return []
    metrics = [m for m in (normalize_metric(i) for i in items if isinstance(i, dict)) if m]
    return metrics


# ── Document generation helpers ──────────────────────────────────────────────
# Where generated reports (PDFs) are written. They live alongside uploads so the
# existing /download endpoint can serve them and they can be re-indexed.
GENERATED_DIR = UPLOAD_DIR
MAX_DOC_CONTEXT_CHARS = 80_000

# Bundled Unicode font. fpdf2's core fonts (Helvetica) are Latin-1 only and choke
# on characters the LLM emits (en-dash, curly quotes, bullets, …), so we register
# DejaVuSans, which covers full Unicode.
FONT_DIR = Path(__file__).parent / "fonts"
_DEJAVU_STYLES = {
    "": "DejaVuSans.ttf",
    "B": "DejaVuSans-Bold.ttf",
    "I": "DejaVuSans-Oblique.ttf",
    "BI": "DejaVuSans-BoldOblique.ttf",
}

# Fallback transliteration used only when the Unicode font is unavailable.
_LATIN1_MAP = {
    "–": "-", "—": "--", "‘": "'", "’": "'",
    "“": '"', "”": '"', "•": "-", "…": "...",
    " ": " ", "→": "->", "−": "-",
}

DOC_PROMPT = ChatPromptTemplate.from_messages([
    ("system",
     "You are a professional report writer for SNAP AI. Using ONLY the source "
     "material provided, write a well-structured document that fulfils the "
     "user's request. Output GitHub-flavored Markdown only — no code fences "
     "around the whole thing, no commentary before or after.\n\n"
     "Requirements:\n"
     "- Start with a single H1 title line ('# Title').\n"
     "- Use H2/H3 headings, bullet lists, and Markdown tables where helpful.\n"
     "- Be accurate to the source; do not invent facts or figures.\n"
     "- Use a clear, professional tone suitable for a business document.\n\n"
     "Source material:\n{context}"),
    ("human", "{instruction}"),
])


def build_doc_context(source: str | None) -> str:
    """Gather the source text for report generation: the whole focused document,
    or the most relevant chunks across all documents when none is focused."""
    if source:
        docs = get_chunks_for_source(source)
    else:
        results = vectorstore._collection.get(include=["documents", "metadatas"])
        documents = results.get("documents") or []
        metadatas = results.get("metadatas") or []
        docs = [
            Document(page_content=t or "", metadata=m or {})
            for t, m in zip(documents, metadatas)
        ]
    return format_docs(docs)[:MAX_DOC_CONTEXT_CHARS]


def extract_title(md_text: str) -> str:
    """Use the first Markdown H1 as the document title, else a default."""
    for line in md_text.splitlines():
        if line.strip().startswith("# "):
            return line.strip()[2:].strip()
    return "Generated Report"


def slugify(text: str) -> str:
    slug = re.sub(r"[^a-z0-9]+", "-", text.lower()).strip("-")
    return (slug or "report")[:50]


def register_pdf_font(pdf: FPDF) -> str:
    """Register the bundled DejaVu family (all styles) so Unicode renders.
    Returns the font family to use, or 'Helvetica' if the fonts are missing."""
    paths = {style: FONT_DIR / name for style, name in _DEJAVU_STYLES.items()}
    if not all(p.is_file() for p in paths.values()):
        return "Helvetica"
    for style, path in paths.items():
        pdf.add_font("DejaVu", style, str(path))
    return "DejaVu"


def to_latin1(text: str) -> str:
    """Best-effort downgrade to Latin-1 for the no-Unicode-font fallback."""
    for uni, ascii_ in _LATIN1_MAP.items():
        text = text.replace(uni, ascii_)
    return text.encode("latin-1", "replace").decode("latin-1")


def markdown_to_pdf(md_text: str, out_path: Path) -> None:
    """Render Markdown to a PDF via HTML (fpdf2 supports headings/lists/tables)."""
    html = markdown.markdown(md_text, extensions=["tables", "sane_lists"])
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    family = register_pdf_font(pdf)
    if family == "Helvetica":
        html = to_latin1(html)  # core font can't render Unicode punctuation
    pdf.set_font(family, size=11)
    pdf.write_html(html)
    pdf.output(str(out_path))


async def run_document_generation(instruction: str, source: str | None) -> dict:
    context = build_doc_context(source)
    if not context.strip():
        raise HTTPException(status_code=400, detail="No document content to work from.")

    chain = DOC_PROMPT | llm | StrOutputParser()
    md_text = await chain.ainvoke({"context": context, "instruction": instruction})

    title = extract_title(md_text)
    filename = f"{slugify(title)}-{datetime.now():%Y%m%d-%H%M%S}.pdf"
    out_path = (GENERATED_DIR / filename).resolve()
    try:
        markdown_to_pdf(md_text, out_path)
    except Exception as exc:
        raise HTTPException(status_code=500, detail=f"PDF generation failed: {exc}")

    return {
        "filename": filename,
        "title": title,
        "markdown": md_text,
        "source": source,
    }


def to_lc_messages(history: list[dict] | None) -> list:
    """Convert [{'role': 'user'|'assistant', 'content': str}, ...] to LangChain messages."""
    msgs = []
    for m in history or []:
        role = m.get("role")
        content = m.get("content", "")
        if not content:
            continue
        if role == "user":
            msgs.append(HumanMessage(content=content))
        elif role == "assistant":
            msgs.append(AIMessage(content=content))
    return msgs


async def run_rag_chain(
    question: str,
    session_id: str,
    source: str | None = None,
    history: list[dict] | None = None,
):
    # If the caller supplies history (persisted in the DB), use it so memory
    # survives restarts and resuming old threads. Otherwise fall back to the
    # in-memory per-session history.
    if history is not None:
        chat_history = to_lc_messages(history)[-20:]
    else:
        chat_history = get_history(session_id)

    # Step 1 — condense follow-up into standalone question if there's history
    if chat_history:
        condense_chain = CONDENSE_PROMPT | llm | StrOutputParser()
        standalone = await condense_chain.ainvoke({"chat_history": chat_history, "question": question})
    else:
        standalone = question

    # Step 2 — retrieve context
    if source:
        # Chat is focused on a specific uploaded document.
        if wants_full_document(question):
            # Summary / overview style request → feed the entire document.
            docs = get_chunks_for_source(source)
        else:
            # Specific question → similarity search, but scoped to this document.
            retriever = vectorstore.as_retriever(
                search_type="similarity",
                search_kwargs={"k": 8, "filter": {"source": source}},
            )
            docs = await retriever.ainvoke(standalone)
            if not docs:  # fall back to the full document if nothing matched
                docs = get_chunks_for_source(source)
    else:
        retriever = vectorstore.as_retriever(search_type="similarity", search_kwargs={"k": 5})
        docs = await retriever.ainvoke(standalone)

    # Step 3 — generate answer (cap context length for very large documents)
    context = format_docs(docs)
    if len(context) > MAX_CONTEXT_CHARS:
        context = context[:MAX_CONTEXT_CHARS]

    rag_chain = RAG_PROMPT | llm | StrOutputParser()
    answer = await rag_chain.ainvoke({
        "context": context,
        "chat_history": chat_history,
        "question": question,
    })

    # Step 4 — update in-memory history only when not using caller-supplied history.
    if history is None:
        mem = get_history(session_id)
        mem.append(HumanMessage(content=question))
        mem.append(AIMessage(content=answer))
        if len(mem) > 12:
            session_histories[session_id] = mem[-12:]

    sources = list({d.metadata.get("source", "unknown") for d in docs})
    return answer, sources


async def run_plain_chain(question: str, session_id: str, history: list[dict] | None = None):
    if history is not None:
        chat_history = to_lc_messages(history)[-20:]
    else:
        chat_history = get_history(session_id)

    plain_chain = PLAIN_PROMPT | llm | StrOutputParser()
    answer = await plain_chain.ainvoke({"chat_history": chat_history, "question": question})

    if history is None:
        mem = get_history(session_id)
        mem.append(HumanMessage(content=question))
        mem.append(AIMessage(content=answer))
        if len(mem) > 12:
            session_histories[session_id] = mem[-12:]

    return answer


# ── Routes ─────────────────────────────────────────────────────────────────────
@app.get("/health")
def health():
    return {"status": "ok", "docs_in_store": vectorstore._collection.count()}


class ChatRequest(BaseModel):
    message: str
    session_id: str = "default"
    source: str | None = None  # focus the chat on a specific uploaded document
    history: list[dict] | None = None  # prior messages [{role, content}] for memory


class ChatResponse(BaseModel):
    answer: str
    sources: list[str]
    doc_count: int
    chart: dict | None = None  # chart/table spec when the prompt asked for one
    document: dict | None = None  # generated-document info when one was created


@app.post("/chat", response_model=ChatResponse)
async def chat(req: ChatRequest):
    doc_count = vectorstore._collection.count()

    try:
        if doc_count == 0:
            answer = await run_plain_chain(req.message, req.session_id, req.history)
            return ChatResponse(answer=answer, sources=[], doc_count=0)

        # If the prompt asks for a generated report/document/PDF, build one from
        # the uploaded material and return a downloadable file.
        if wants_document(req.message):
            # Let HTTPExceptions (e.g. no content, PDF failure) surface to the
            # client so problems are visible rather than silently degrading.
            doc = await run_document_generation(req.message, resolve_source(req.source))
            answer = (
                f"I've generated **{doc['title']}**. Download it below, "
                "or add it back to my knowledge base."
            )
            return ChatResponse(
                answer=answer,
                sources=[doc["source"]] if doc.get("source") else [],
                doc_count=doc_count,
                document=doc,
            )

        # If the prompt asks for a chart/graph/table and we can pin down a
        # document, build a chart spec inline. Fall back to a normal answer if
        # the request can't be turned into a chart.
        if wants_chart(req.message):
            # Resolve to a focused/only document, or None → chart across all docs.
            source = resolve_source(req.source)
            try:
                spec = await run_visualization(source, req.message)
                answer = spec.get("notes") or spec.get("title") or "Here's the chart you asked for."
                return ChatResponse(
                    answer=answer,
                    sources=[source] if source else [],
                    doc_count=doc_count,
                    chart=spec,
                )
            except HTTPException:
                pass  # not chartable → continue with a normal RAG answer

        answer, sources = await run_rag_chain(req.message, req.session_id, req.source, req.history)
        return ChatResponse(answer=answer, sources=sources, doc_count=doc_count)
    except ChatGoogleGenerativeAIError as e:
        # Surface Gemini quota/rate-limit errors as a clean 429 instead of a 500
        # crash, so the client can show a friendly "try again later" message.
        msg = str(e)
        if "RESOURCE_EXHAUSTED" in msg or "429" in msg:
            raise HTTPException(
                status_code=429,
                detail=(
                    f"Gemini quota exceeded for model '{GEMINI_MODEL}'. "
                    "You've hit the free-tier request limit — wait for the quota "
                    "to reset, switch GEMINI_MODEL, or enable billing."
                ),
            )
        raise HTTPException(status_code=502, detail=f"Gemini error: {msg}")


class VisualizeRequest(BaseModel):
    instruction: str
    source: str | None = None  # which uploaded document to chart


@app.post("/visualize")
async def visualize(req: VisualizeRequest):
    if vectorstore._collection.count() == 0:
        raise HTTPException(status_code=400, detail="No documents uploaded yet.")

    # Resolve which document to chart: the requested one, or the only one indexed.
    source = resolve_source(req.source)
    if not source:
        raise HTTPException(
            status_code=400,
            detail="Select a document to chart (multiple documents are indexed).",
        )

    try:
        return await run_visualization(source, req.instruction)
    except ChatGoogleGenerativeAIError as e:
        msg = str(e)
        if "RESOURCE_EXHAUSTED" in msg or "429" in msg:
            raise HTTPException(
                status_code=429,
                detail=(
                    f"Gemini quota exceeded for model '{GEMINI_MODEL}'. "
                    "You've hit the free-tier request limit — wait for the quota "
                    "to reset, switch GEMINI_MODEL, or enable billing."
                ),
            )
        raise HTTPException(status_code=502, detail=f"Gemini error: {msg}")


class ExtractMetricsRequest(BaseModel):
    source: str  # which uploaded document to extract dashboard metrics from


@app.post("/extract-metrics")
async def extract_metrics_endpoint(req: ExtractMetricsRequest):
    try:
        metrics = await extract_metrics(req.source)
        return {"source": req.source, "metrics": metrics}
    except ChatGoogleGenerativeAIError as e:
        msg = str(e)
        if "RESOURCE_EXHAUSTED" in msg or "429" in msg:
            raise HTTPException(
                status_code=429,
                detail=(
                    f"Gemini quota exceeded for model '{GEMINI_MODEL}'. "
                    "You've hit the free-tier request limit — wait for the quota "
                    "to reset, switch GEMINI_MODEL, or enable billing."
                ),
            )
        raise HTTPException(status_code=502, detail=f"Gemini error: {msg}")


class GenerateDocRequest(BaseModel):
    instruction: str
    source: str | None = None  # which uploaded document to base the report on


@app.post("/generate-document")
async def generate_document(req: GenerateDocRequest):
    if vectorstore._collection.count() == 0:
        raise HTTPException(status_code=400, detail="No documents uploaded yet.")
    try:
        return await run_document_generation(req.instruction, resolve_source(req.source))
    except ChatGoogleGenerativeAIError as e:
        msg = str(e)
        if "RESOURCE_EXHAUSTED" in msg or "429" in msg:
            raise HTTPException(
                status_code=429,
                detail=(
                    f"Gemini quota exceeded for model '{GEMINI_MODEL}'. "
                    "You've hit the free-tier request limit — wait for the quota "
                    "to reset, switch GEMINI_MODEL, or enable billing."
                ),
            )
        raise HTTPException(status_code=502, detail=f"Gemini error: {msg}")


class IngestRequest(BaseModel):
    filename: str  # a file already present in the uploads/generated directory


@app.post("/ingest")
def ingest_document(req: IngestRequest):
    """Index a file already on disk (e.g. a generated report) into the vector
    store so the AI can answer questions about it."""
    safe_name = Path(req.filename).name
    path = (UPLOAD_DIR / safe_name).resolve()
    if UPLOAD_DIR.resolve() not in path.parents or not path.is_file():
        raise HTTPException(status_code=404, detail=f"File '{safe_name}' not found")
    try:
        chunks_indexed = index_file(path, safe_name)
    except Exception as exc:
        raise HTTPException(status_code=500, detail=str(exc))
    return {
        "filename": safe_name,
        "chunks_indexed": chunks_indexed,
        "total_docs": vectorstore._collection.count(),
    }


@app.post("/upload")
async def upload_document(file: UploadFile = File(...)):
    if not file.filename:
        raise HTTPException(status_code=400, detail="No filename provided.")

    filename = file.filename
    suffix = Path(filename).suffix.lower()
    if suffix not in ALLOWED_EXTENSIONS:
        raise HTTPException(
            status_code=400,
            detail=f"Unsupported type '{suffix}'. Allowed: {list(ALLOWED_EXTENSIONS)}",
        )

    dest = UPLOAD_DIR / filename
    with dest.open("wb") as f:
        shutil.copyfileobj(file.file, f)

    try:
        chunks_indexed = index_file(dest, filename)
    except Exception as exc:
        dest.unlink(missing_ok=True)
        raise HTTPException(status_code=500, detail=str(exc))

    return {
        "filename": filename,
        "chunks_indexed": chunks_indexed,
        "total_docs": vectorstore._collection.count(),
    }


@app.delete("/documents")
def clear_documents():
    global vectorstore
    vectorstore._client.delete_collection("snap_ai_docs")
    vectorstore = Chroma(
        collection_name="snap_ai_docs",
        embedding_function=embeddings,
        persist_directory=str(CHROMA_DIR),
    )
    session_histories.clear()

    # Also remove the original uploaded files and any generated reports on disk
    # so nothing lingers after a clear (GENERATED_DIR == UPLOAD_DIR).
    removed = 0
    for path in UPLOAD_DIR.iterdir():
        if path.is_file():
            try:
                path.unlink()
                removed += 1
            except OSError:
                pass

    return {"status": "cleared", "files_removed": removed}


@app.delete("/session/{session_id}")
def clear_session(session_id: str):
    session_histories.pop(session_id, None)
    return {"status": "cleared", "session_id": session_id}


def _source_mtime(name: str) -> float:
    """Modification time of a source's file on disk (newest uploads have the
    largest mtime). Sources whose file is gone sort last."""
    try:
        return (UPLOAD_DIR / Path(name).name).stat().st_mtime
    except OSError:
        return 0.0


@app.get("/documents")
def list_documents():
    count = vectorstore._collection.count()
    if count == 0:
        return {"documents": [], "total_chunks": 0}
    results = vectorstore._collection.get(include=["metadatas"])
    metadatas = results.get("metadatas") or []
    sources = {str(m.get("source", "unknown")) for m in metadatas}
    # Newest uploaded first.
    ordered = sorted(sources, key=_source_mtime, reverse=True)
    return {"documents": ordered, "total_chunks": count}


@app.get("/download/{filename}")
def download_document(filename: str):
    # Resolve against UPLOAD_DIR and reject anything that escapes it (path traversal).
    safe_name = Path(filename).name
    file_path = (UPLOAD_DIR / safe_name).resolve()
    if UPLOAD_DIR.resolve() not in file_path.parents or not file_path.is_file():
        raise HTTPException(status_code=404, detail=f"File '{safe_name}' not found")
    return FileResponse(
        path=str(file_path),
        filename=safe_name,
        media_type="application/octet-stream",
    )
