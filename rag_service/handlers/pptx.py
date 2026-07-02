from pathlib import Path

from .base import Parsed


def parse(path: Path) -> Parsed:
    from pptx import Presentation

    prs = Presentation(str(path))
    texts = []
    for slide in prs.slides:
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:                       # note: text_frame, not textframe
                for paragraph in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in paragraph.runs).strip()
                    if line:
                        parts.append(line)
            elif shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells if c.text.strip()]
                    if cells:
                        parts.append(" | ".join(cells))
        if parts:
            texts.append("\n".join(parts))
    return Parsed(text_chunks=texts)
